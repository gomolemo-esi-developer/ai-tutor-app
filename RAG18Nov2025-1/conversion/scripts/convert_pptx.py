import sys
from pathlib import Path

sys.path.append(str(Path(__file__).parent.parent))

from utils.progress_tracker import ProgressTracker
from utils.file_handler import FileHandler

class PPTXConverter:
    def __init__(self, output_dir):
        self.file_handler = FileHandler(output_dir)
        
    def extract_text_from_slide(self, slide):
        text_parts = []
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text_parts.append(shape.text)
        
        if hasattr(slide, "notes_slide") and slide.notes_slide:
            notes_shape = slide.notes_slide.notes_text_frame
            if notes_shape and notes_shape.text:
                text_parts.append(f"\n[Notes: {notes_shape.text}]")
        
        return "\n".join(text_parts)
    
    def convert(self, pptx_path):
        pptx_path = self.file_handler.validate_input_file(
            pptx_path, 
            allowed_extensions=['.pptx', '.ppt']
        )
        
        print(f"\nConverting PowerPoint: {pptx_path.name}")
        print(f"Size: {self.file_handler.get_file_size_mb(pptx_path):.2f} MB\n")
        
        try:
            from pptx import Presentation
            
            presentation = Presentation(str(pptx_path))
            total_slides = len(presentation.slides)
            
            print(f"Total slides: {total_slides}\n")
            
            progress = ProgressTracker("Extracting text", "slide")
            pbar = progress.create_bar(total_slides)
            
            extracted_text_parts = []
            
            for i, slide in enumerate(presentation.slides, 1):
                pbar.set_description(f"Processing slide {i}/{total_slides}")
                
                slide_text = self.extract_text_from_slide(slide)
                
                if slide_text.strip():
                    extracted_text_parts.append(f"=== Slide {i} ===\n{slide_text}\n")
                else:
                    extracted_text_parts.append(f"=== Slide {i} ===\n[No text content]\n")
                
                pbar.update(1)
            
            pbar.close()
            
            full_text = "\n".join(extracted_text_parts)
            
            output_path = self.file_handler.get_output_path(
                pptx_path,
                suffix="_slides",
                extension=".txt"
            )
            
            self.file_handler.save_text(full_text, output_path)
            
            print(f"\n✅ Extracted text saved to: {output_path}")
            return str(output_path)
            
        except ImportError:
            raise ImportError("python-pptx not installed. Install with: pip install python-pptx")
        except Exception as e:
            raise RuntimeError(f"PPTX conversion failed: {str(e)}")

def main():
    if len(sys.argv) < 2:
        print("Usage: python convert_pptx.py <pptx_file>")
        sys.exit(1)
    
    pptx_file = sys.argv[1]
    
    script_dir = Path(__file__).parent.parent
    output_dir = script_dir / "output"
    
    converter = PPTXConverter(output_dir=output_dir)
    
    try:
        converter.convert(pptx_file)
    except Exception as e:
        print(f"\n❌ Error: {str(e)}")
        sys.exit(1)

if __name__ == "__main__":
    main()

