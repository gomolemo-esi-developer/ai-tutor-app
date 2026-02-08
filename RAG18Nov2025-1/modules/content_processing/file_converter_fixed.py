from pathlib import Path
from typing import Optional, Callable, Dict
import sys
import tempfile
import os

sys.path.append(str(Path(__file__).parent.parent.parent))
from modules.shared.logger import setup_logger

logger = setup_logger(__name__)

class FileConverter:
    def __init__(self, models_dir: str, output_dir: str):
        self.models_dir = Path(models_dir)
        self.output_dir = Path(output_dir)
        self.models_dir.mkdir(parents=True, exist_ok=True)
        self.output_dir.mkdir(parents=True, exist_ok=True)
        
        self.audio_extensions = {'.mp3', '.wav', '.m4a', '.flac', '.ogg'}
        self.video_extensions = {'.mp4', '.avi', '.mov', '.mkv', '.webm'}
        self.pdf_extensions = {'.pdf'}
        self.pptx_extensions = {'.pptx', '.ppt'}
        self.code_extensions = {
            '.cs', '.py', '.java', '.cpp', '.c', '.h', '.hpp',
            '.js', '.ts', '.jsx', '.tsx', '.go', '.rs', '.rb',
            '.php', '.swift', '.kt', '.scala', '.r', '.m', '.mm'
        }
        self.image_extensions = {'.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.tif', '.webp'}
        self.text_extensions = {'.txt', '.md', '.csv'}
        
    def get_file_type(self, file_path: Path) -> str:
        ext = file_path.suffix.lower()
        if ext in self.audio_extensions:
            return 'audio'
        elif ext in self.video_extensions:
            return 'video'
        elif ext in self.pdf_extensions:
            return 'pdf'
        elif ext in self.pptx_extensions:
            return 'pptx'
        elif ext in self.code_extensions:
            return 'code'
        elif ext in self.image_extensions:
            return 'image'
        elif ext in self.text_extensions:
            return 'text'
        else:
            return 'unknown'
    
    def convert_to_text(self, file_path: str, progress_callback: Optional[Callable] = None) -> Optional[str]:
        file_path = Path(file_path)
        
        if not file_path.exists():
            logger.error(f"File not found: {file_path}")
            return None
        
        file_type = self.get_file_type(file_path)
        logger.info(f"Converting {file_type} file: {file_path.name}")
        
        if progress_callback:
            progress_callback(f"Detected file type: {file_type}")
        
        try:
            if file_type == 'audio':
                return self._convert_audio(file_path, progress_callback)
            elif file_type == 'video':
                return self._convert_video(file_path, progress_callback)
            elif file_type == 'pdf':
                return self._convert_pdf(file_path, progress_callback)
            elif file_type == 'pptx':
                return self._convert_pptx(file_path, progress_callback)
            elif file_type == 'code':
                return self._convert_code(file_path, progress_callback)
            elif file_type == 'image':
                return self._convert_image(file_path, progress_callback)
            elif file_type == 'text':
                return self._convert_text(file_path, progress_callback)
            else:
                logger.error(f"Unsupported file type: {file_path.suffix}")
                return None
                
        except Exception as e:
            logger.error(f"Conversion failed for {file_path.name}: {str(e)}")
            raise
    
    def _convert_audio(self, file_path: Path, callback: Optional[Callable] = None) -> str:
        if callback:
            callback("Loading Whisper transcription model...")
        
        try:
            from faster_whisper import WhisperModel
            
            model = WhisperModel(
                "base",
                device="cpu",
                compute_type="int8",
                download_root=str(self.models_dir)
            )
            
            if callback:
                callback(f"Transcribing audio: {file_path.name}")
            
            segments, info = model.transcribe(str(file_path), beam_size=5)
            
            transcription_parts = []
            for segment in segments:
                transcription_parts.append(segment.text)
            
            text = " ".join(transcription_parts)
            logger.info(f"✅ Transcribed {file_path.name} ({len(text)} chars)")
            
            return text
            
        except ImportError:
            logger.error("faster-whisper not installed")
            raise ImportError("faster-whisper required. Install: pip install faster-whisper")
    
    def _convert_video(self, file_path: Path, callback: Optional[Callable] = None) -> str:
        import subprocess
        
        if callback:
            callback(f"Extracting audio from video: {file_path.name}")
        
        temp_audio = None
        try:
            with tempfile.NamedTemporaryFile(suffix='.mp3', delete=False) as tmp:
                temp_audio = tmp.name
            
            cmd = [
                'ffmpeg',
                '-i', str(file_path),
                '-vn',
                '-acodec', 'libmp3lame',
                '-q:a', '2',
                '-y',
                str(temp_audio)
            ]
            
            result = subprocess.run(
                cmd,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                text=True
            )
            
            if result.returncode != 0:
                raise RuntimeError(f"ffmpeg failed: {result.stderr}")
            
            if callback:
                callback("Audio extracted, starting transcription...")
            
            text = self._convert_audio(Path(temp_audio), callback)
            
            return text
            
        finally:
            if temp_audio and os.path.exists(temp_audio):
                os.remove(temp_audio)
    
    def _convert_pdf(self, file_path: Path, callback: Optional[Callable] = None) -> str:
        """Convert PDF to text using PyPDF2 instead of pdf2image + OCR"""
        if callback:
            callback(f"Extracting text from PDF...")
        
        try:
            # First try: use PyPDF2 for direct text extraction (fastest)
            try:
                import PyPDF2
                
                extracted_text = []
                with open(file_path, 'rb') as f:
                    pdf_reader = PyPDF2.PdfReader(f)
                    num_pages = len(pdf_reader.pages)
                    
                    for i, page in enumerate(pdf_reader.pages, 1):
                        if callback and i % 5 == 0:
                            callback(f"Extracting text from page {i}/{num_pages}")
                        
                        page_text = page.extract_text()
                        if page_text:
                            extracted_text.append(f"--- Page {i} ---\n{page_text}")
                
                if extracted_text:
                    text = "\n\n".join(extracted_text)
                    logger.info(f"✅ Extracted {num_pages} pages from {file_path.name} using PyPDF2")
                    return text
                else:
                    # If PyPDF2 failed, fall through to OCR
                    logger.warning(f"PyPDF2 extracted no text from {file_path.name}, trying OCR...")
            
            except ImportError:
                logger.info("PyPDF2 not installed, using OCR fallback...")
            
            # Fallback: use OCR for scanned PDFs
            from pdf2image import convert_from_path
            import numpy as np
            
            images = convert_from_path(str(file_path))
            
            if callback:
                callback(f"Processing {len(images)} pages with OCR...")
            
            try:
                import easyocr
                reader = easyocr.Reader(['en'], gpu=False, verbose=False)
                
                extracted_pages = []
                for i, image in enumerate(images, 1):
                    if callback:
                        callback(f"OCR page {i}/{len(images)}")
                    
                    # Convert PIL Image to numpy array for easyocr
                    image_np = np.array(image)
                    result = reader.readtext(image_np, detail=0, paragraph=False)
                    page_text = "\n".join(result) if result else ""
                    extracted_pages.append(f"--- Page {i} ---\n{page_text}")
                
                text = "\n\n".join(extracted_pages)
                logger.info(f"✅ Extracted {len(images)} pages from {file_path.name} using OCR")
                
                return text
                
            except ImportError:
                logger.error("easyocr not installed")
                raise ImportError("easyocr required. Install: pip install easyocr")
                
        except ImportError as e:
            logger.error(f"pdf2image not installed: {str(e)}")
            raise ImportError("pdf2image required. Install: pip install pdf2image")
    
    def _convert_pptx(self, file_path: Path, callback: Optional[Callable] = None) -> str:
        if callback:
            callback(f"Extracting text from PowerPoint...")
        
        try:
            from pptx import Presentation
            
            presentation = Presentation(str(file_path))
            total_slides = len(presentation.slides)
            
            extracted_slides = []
            for i, slide in enumerate(presentation.slides, 1):
                if callback and i % 5 == 0:
                    callback(f"Processing slide {i}/{total_slides}")
                
                slide_text_parts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text_parts.append(shape.text)
                
                slide_text = "\n".join(slide_text_parts)
                if slide_text.strip():
                    extracted_slides.append(f"=== Slide {i} ===\n{slide_text}")
            
            text = "\n\n".join(extracted_slides)
            logger.info(f"✅ Extracted {total_slides} slides from {file_path.name}")
            
            return text
            
        except ImportError:
            logger.error("python-pptx not installed")
            raise ImportError("python-pptx required. Install: pip install python-pptx")
    
    def _convert_code(self, file_path: Path, callback: Optional[Callable] = None) -> str:
        if callback:
            callback(f"Reading code file...")
        
        encodings = ['utf-8', 'utf-16', 'latin-1', 'cp1252']
        
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    content = f.read()
                
                metadata = f"File: {file_path.name}\n"
                metadata += f"Type: {file_path.suffix}\n"
                metadata += f"Lines: {len(content.splitlines())}\n"
                metadata += "=" * 50 + "\n\n"
                
                text = metadata + content
                logger.info(f"✅ Read {file_path.name} ({len(content)} chars)")
                
                return text
            except (UnicodeDecodeError, UnicodeError):
                continue
        
        raise RuntimeError(f"Unable to decode {file_path.name}")
    
    def _convert_image(self, file_path: Path, callback: Optional[Callable] = None) -> str:
        if callback:
            callback(f"Running OCR on image...")
        
        try:
            import easyocr
            import numpy as np
            from PIL import Image
            
            # Load image and convert to numpy array
            img = Image.open(file_path)
            img_np = np.array(img)
            
            reader = easyocr.Reader(['en'], gpu=False, verbose=False)
            result = reader.readtext(img_np, detail=0, paragraph=False)
            
            if not result:
                return "No text detected in image"
            
            text = "\n".join(result)
            logger.info(f"✅ Extracted {len(result)} lines from {file_path.name}")
            
            return text
            
        except ImportError:
            logger.error("easyocr not installed")
            raise ImportError("easyocr required. Install: pip install easyocr")
    
    def _convert_text(self, file_path: Path, callback: Optional[Callable] = None) -> str:
        if callback:
            callback(f"Reading text file...")
        
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                text = f.read()
            
            logger.info(f"✅ Read {file_path.name} ({len(text)} chars)")
            return text
            
        except Exception as e:
            logger.error(f"Failed to read {file_path.name}: {str(e)}")
            raise
    
    def get_supported_extensions(self) -> Dict[str, list]:
        return {
            'audio': list(self.audio_extensions),
            'video': list(self.video_extensions),
            'pdf': list(self.pdf_extensions),
            'pptx': list(self.pptx_extensions),
            'code': list(self.code_extensions),
            'image': list(self.image_extensions),
            'text': list(self.text_extensions)
        }
    
    def is_supported(self, file_path: str) -> bool:
        ext = Path(file_path).suffix.lower()
        all_extensions = (
            self.audio_extensions | self.video_extensions | self.pdf_extensions |
            self.pptx_extensions | self.code_extensions | self.image_extensions |
            self.text_extensions
        )
        return ext in all_extensions
