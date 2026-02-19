from pathlib import Path
from typing import Optional, Callable, Dict
import sys
import tempfile
import os
import subprocess

sys.path.append(str(Path(__file__).parent.parent.parent))
from modules.shared.logger import setup_logger

logger = setup_logger(__name__)

class FileConverter:
    def __init__(self, models_dir: str, output_dir: str):
        self.models_dir = Path(models_dir)
        self.output_dir = Path(output_dir)
        self.models_dir.mkdir(parents=True, exist_ok=True)
        self.output_dir.mkdir(parents=True, exist_ok=True)
        
        self.audio_extensions = {'.mp3', '.wav', '.m4a', '.flac', '.ogg'}
        self.video_extensions = {'.mp4', '.avi', '.mov', '.mkv', '.webm'}
        self.pdf_extensions = {'.pdf'}
        self.pptx_extensions = {'.pptx', '.ppt'}
        self.code_extensions = {
            '.cs', '.py', '.java', '.cpp', '.c', '.h', '.hpp',
            '.js', '.ts', '.jsx', '.tsx', '.go', '.rs', '.rb',
            '.php', '.swift', '.kt', '.scala', '.r', '.m', '.mm'
        }
        self.image_extensions = {'.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.tif', '.webp'}
        self.text_extensions = {'.txt', '.md', '.csv'}
        self.docx_extensions = {'.docx', '.doc'}
        self.excel_extensions = {'.xlsx', '.xls', '.csv'}
        
    def get_file_type(self, file_path: Path) -> str:
        ext = file_path.suffix.lower()
        if ext in self.audio_extensions:
            return 'audio'
        elif ext in self.video_extensions:
            return 'video'
        elif ext in self.pdf_extensions:
            return 'pdf'
        elif ext in self.pptx_extensions:
            return 'pptx'
        elif ext in self.docx_extensions:
            return 'docx'
        elif ext in self.excel_extensions:
            return 'excel'
        elif ext in self.code_extensions:
            return 'code'
        elif ext in self.image_extensions:
            return 'image'
        elif ext in self.text_extensions:
            return 'text'
        else:
            return 'unknown'
    
    def convert_to_text(self, file_path: str, progress_callback: Optional[Callable] = None) -> Optional[str]:
        file_path = Path(file_path)
        
        if not file_path.exists():
            logger.error(f"File not found: {file_path}")
            return None
        
        file_type = self.get_file_type(file_path)
        logger.info(f"Converting {file_type} file: {file_path.name}")
        
        if progress_callback:
            progress_callback(f"Detected file type: {file_type}")
        
        try:
            if file_type == 'audio':
                return self._convert_audio(file_path, progress_callback)
            elif file_type == 'video':
                return self._convert_video(file_path, progress_callback)
            elif file_type == 'pdf':
                return self._convert_pdf(file_path, progress_callback)
            elif file_type == 'pptx':
                return self._convert_pptx(file_path, progress_callback)
            elif file_type == 'docx':
                return self._convert_docx(file_path, progress_callback)
            elif file_type == 'excel':
                return self._convert_excel(file_path, progress_callback)
            elif file_type == 'code':
                return self._convert_code(file_path, progress_callback)
            elif file_type == 'image':
                return self._convert_image(file_path, progress_callback)
            elif file_type == 'text':
                return self._convert_text(file_path, progress_callback)
            else:
                logger.error(f"Unsupported file type: {file_path.suffix}")
                return None
                
        except Exception as e:
            logger.error(f"Conversion failed for {file_path.name}: {str(e)}")
            raise
    
    def _convert_audio(self, file_path: Path, callback: Optional[Callable] = None) -> str:
        if callback:
            callback("Loading Whisper transcription model...")
        
        try:
            from faster_whisper import WhisperModel
            import gc
            
            # FIX (2026-02-19): Reduce memory footprint for large audio files
            # - Use tiny model instead of base for 37MB+ files to fit in 512MB memory
            # - Use beam_size=1 (greedy) instead of 5 to reduce memory
            # - Force CPU-only with int8 quantization
            # - Aggressive garbage collection
            model = WhisperModel(
                "tiny",  # CHANGED: Use tiny model (139M) instead of base (140M with less compression)
                device="cpu",
                compute_type="int8",  # Quantized to int8 for memory efficiency
                download_root=str(self.models_dir)
            )
            
            if callback:
                callback(f"Transcribing audio: {file_path.name}")
            
            # FIX: Use greedy decoding (beam_size=1) instead of beam search to save memory
            segments, info = model.transcribe(
                str(file_path), 
                beam_size=1,  # CHANGED: Greedy decoding uses much less memory
                language="en",
                vad_filter=True,  # Voice activity detection to skip silence
                vad_parameters=dict(min_speech_duration_ms=250)
            )
            
            transcription_parts = []
            for segment in segments:
                transcription_parts.append(segment.text)
                # Clear segment to free memory
                del segment
            
            text = " ".join(transcription_parts)
            
            # Force garbage collection after transcription
            del model
            gc.collect()
            
            logger.info(f"✅ Transcribed {file_path.name} ({len(text)} chars)")
            
            return text
            
        except ImportError:
            logger.error("faster-whisper not installed")
            raise ImportError("faster-whisper required. Install: pip install faster-whisper")
    
    def _convert_video(self, file_path: Path, callback: Optional[Callable] = None) -> str:
        import subprocess
        
        if callback:
            callback(f"Extracting audio from video: {file_path.name}")
        
        temp_audio = None
        try:
            with tempfile.NamedTemporaryFile(suffix='.mp3', delete=False) as tmp:
                temp_audio = tmp.name
            
            cmd = [
                'ffmpeg',
                '-i', str(file_path),
                '-vn',
                '-acodec', 'libmp3lame',
                '-q:a', '2',
                '-y',
                str(temp_audio)
            ]
            
            result = subprocess.run(
                cmd,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                text=True
            )
            
            if result.returncode != 0:
                raise RuntimeError(f"ffmpeg failed: {result.stderr}")
            
            if callback:
                callback("Audio extracted, starting transcription...")
            
            text = self._convert_audio(Path(temp_audio), callback)
            
            return text
            
        finally:
            if temp_audio and os.path.exists(temp_audio):
                os.remove(temp_audio)
    
    def _convert_pdf(self, file_path: Path, callback: Optional[Callable] = None) -> str:
        """Convert PDF to text using PyPDF2 first, then OCR as fallback"""
        if callback:
            callback(f"Extracting text from PDF...")
        
        try:
            # First try: use PyPDF2 for direct text extraction (fastest)
            try:
                import PyPDF2
                
                extracted_text = []
                with open(file_path, 'rb') as f:
                    pdf_reader = PyPDF2.PdfReader(f)
                    num_pages = len(pdf_reader.pages)
                    
                    for i, page in enumerate(pdf_reader.pages, 1):
                        if callback and i % 5 == 0:
                            callback(f"Extracting text from page {i}/{num_pages}")
                        
                        page_text = page.extract_text()
                        if page_text:
                            extracted_text.append(f"--- Page {i} ---\n{page_text}")
                
                if extracted_text:
                    text = "\n\n".join(extracted_text)
                    logger.info(f"✅ Extracted {num_pages} pages from {file_path.name} using PyPDF2")
                    return text
                else:
                    # If PyPDF2 failed, fall through to OCR
                    logger.warning(f"PyPDF2 extracted no text from {file_path.name}, trying OCR...")
            
            except ImportError:
                logger.info("PyPDF2 not installed, using OCR fallback...")
            
            # Fallback: use OCR for scanned PDFs
            from pdf2image import convert_from_path
            import numpy as np
            
            if callback:
                callback("Converting PDF pages to images...")
            
            images = convert_from_path(str(file_path))
            
            if callback:
                callback(f"Processing {len(images)} pages with OCR...")
            
            try:
                import easyocr
                reader = easyocr.Reader(['en'], gpu=False, verbose=False)
                
                extracted_pages = []
                for i, image in enumerate(images, 1):
                    if callback:
                        callback(f"OCR page {i}/{len(images)}")
                    
                    # Convert PIL Image to numpy array for easyocr
                    image_np = np.array(image)
                    result = reader.readtext(image_np, detail=0, paragraph=False)
                    page_text = "\n".join(result) if result else ""
                    extracted_pages.append(f"--- Page {i} ---\n{page_text}")
                
                text = "\n\n".join(extracted_pages)
                logger.info(f"✅ Extracted {len(images)} pages from {file_path.name} using OCR")
                
                return text
                
            except ImportError:
                logger.error("easyocr not installed")
                raise ImportError("easyocr required. Install: pip install easyocr")
                
        except ImportError as e:
            logger.error(f"pdf2image not installed: {str(e)}")
            raise ImportError("pdf2image required. Install: pip install pdf2image")
    
    def _convert_pptx(self, file_path: Path, callback: Optional[Callable] = None) -> str:
        if callback:
            callback(f"Extracting text from PowerPoint...")
        
        try:
            from pptx import Presentation
            
            # FIX (2026-02-19): Handle both .pptx and .ppt files
            # python-pptx only works with .pptx format
            # For .ppt files (old format), try converting to .pptx first using LibreOffice
            if file_path.suffix.lower() == '.ppt':
                if callback:
                    callback("Converting .ppt to .pptx format...")
                
                try:
                    import subprocess
                    import tempfile
                    
                    # Use LibreOffice to convert .ppt to .pptx
                    with tempfile.TemporaryDirectory() as tmpdir:
                        converted_path = Path(tmpdir) / f"{file_path.stem}.pptx"
                        
                        cmd = [
                            'libreoffice',
                            '--headless',
                            '--norestore',
                            '--invisible',
                            '--nofirststartwizard',
                            '--convert-to', 'pptx',
                            '--outdir', str(tmpdir),
                            str(file_path)
                        ]
                        
                        result = subprocess.run(cmd, capture_output=True, text=True, timeout=60)
                        
                        if result.returncode != 0:
                            # LibreOffice conversion failed, log details
                            logger.warning(f"LibreOffice conversion failed: {result.stderr}, using fallback...")
                            return self._extract_ppt_fallback(file_path, callback)
                        
                        if not converted_path.exists():
                            logger.warning(f"Conversion succeeded but file not found at {converted_path}")
                            return self._extract_ppt_fallback(file_path, callback)
                        
                        # Process converted .pptx
                        presentation = Presentation(str(converted_path))
                        file_to_process = converted_path
                except Exception as e:
                    logger.warning(f"LibreOffice conversion failed: {str(e)}, using fallback...")
                    return self._extract_ppt_fallback(file_path, callback)
            else:
                # For .pptx files, open directly
                presentation = Presentation(str(file_path))
                file_to_process = file_path
            
            total_slides = len(presentation.slides)
            
            extracted_slides = []
            for i, slide in enumerate(presentation.slides, 1):
                if callback and i % 5 == 0:
                    callback(f"Processing slide {i}/{total_slides}")
                
                slide_text_parts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text_parts.append(shape.text)
                
                slide_text = "\n".join(slide_text_parts)
                if slide_text.strip():
                    extracted_slides.append(f"=== Slide {i} ===\n{slide_text}")
            
            text = "\n\n".join(extracted_slides)
            logger.info(f"✅ Extracted {total_slides} slides from {file_path.name}")
            
            return text
            
        except ImportError:
            logger.error("python-pptx not installed")
            raise ImportError("python-pptx required. Install: pip install python-pptx")
    
    def _extract_ppt_fallback(self, file_path: Path, callback: Optional[Callable] = None) -> str:
        """Fallback method for .ppt files - returns error message"""
        logger.error(f"Cannot process .ppt file (legacy format). File: {file_path.name}")
        return f"[Presentation file: {file_path.name}] - Cannot extract text from .ppt (legacy PowerPoint format). Please convert your file to .pptx format and upload again.\n\nHow to convert:\n1. Open the .ppt file in Microsoft PowerPoint or LibreOffice Impress\n2. Go to File > Save As\n3. Change the format to 'PowerPoint Presentation (.pptx)'\n4. Upload the new .pptx file"
    
    def _convert_docx(self, file_path: Path, callback: Optional[Callable] = None) -> str:
        """Extract text from Word documents"""
        if callback:
            callback(f"Extracting text from Word document...")
        
        try:
            from docx import Document
            
            doc = Document(str(file_path))
            total_paragraphs = len(doc.paragraphs)
            
            extracted_text = []
            
            # Extract paragraph text
            for i, para in enumerate(doc.paragraphs, 1):
                if callback and i % 10 == 0:
                    callback(f"Processing paragraph {i}/{total_paragraphs}")
                
                if para.text.strip():
                    extracted_text.append(para.text)
            
            # Extract table data if present
            if doc.tables:
                extracted_text.append("\n\n=== TABLES ===\n")
                for table_idx, table in enumerate(doc.tables, 1):
                    extracted_text.append(f"\n--- Table {table_idx} ---")
                    for row in table.rows:
                        row_cells = [cell.text.strip() for cell in row.cells]
                        extracted_text.append(" | ".join(row_cells))
            
            text = "\n".join(extracted_text)
            logger.info(f"✅ Extracted {total_paragraphs} paragraphs from {file_path.name}")
            
            return text
            
        except ImportError:
            logger.error("python-docx not installed")
            raise ImportError("python-docx required. Install: pip install python-docx")
    
    def _convert_excel(self, file_path: Path, callback: Optional[Callable] = None) -> str:
        """Extract text from Excel files"""
        if callback:
            callback(f"Extracting data from Excel file...")
        
        try:
            import openpyxl
            
            workbook = openpyxl.load_workbook(str(file_path))
            extracted_text = []
            
            for sheet_idx, sheet_name in enumerate(workbook.sheetnames, 1):
                if callback:
                    callback(f"Processing sheet {sheet_idx}/{len(workbook.sheetnames)}: {sheet_name}")
                
                sheet = workbook[sheet_name]
                extracted_text.append(f"\n=== Sheet: {sheet_name} ===\n")
                
                # Extract all cell values
                for row in sheet.iter_rows(values_only=True):
                    row_values = []
                    for cell in row:
                        if cell is not None:
                            row_values.append(str(cell).strip())
                    
                    if row_values and any(row_values):  # Only add non-empty rows
                        extracted_text.append(" | ".join(row_values))
            
            text = "\n".join(extracted_text)
            logger.info(f"✅ Extracted {len(workbook.sheetnames)} sheets from {file_path.name}")
            
            return text
            
        except ImportError:
            logger.error("openpyxl not installed")
            raise ImportError("openpyxl required. Install: pip install openpyxl")
    
    def _convert_code(self, file_path: Path, callback: Optional[Callable] = None) -> str:
        if callback:
            callback(f"Reading code file...")
        
        encodings = ['utf-8', 'utf-16', 'latin-1', 'cp1252']
        
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    content = f.read()
                
                metadata = f"File: {file_path.name}\n"
                metadata += f"Type: {file_path.suffix}\n"
                metadata += f"Lines: {len(content.splitlines())}\n"
                metadata += "=" * 50 + "\n\n"
                
                text = metadata + content
                logger.info(f"✅ Read {file_path.name} ({len(content)} chars)")
                
                return text
            except (UnicodeDecodeError, UnicodeError):
                continue
        
        raise RuntimeError(f"Unable to decode {file_path.name}")
    
    def _convert_image(self, file_path: Path, callback: Optional[Callable] = None) -> str:
        if callback:
            callback(f"Running OCR on image...")
        
        try:
            import easyocr
            import numpy as np
            from PIL import Image
            
            # Load image and convert to numpy array
            img = Image.open(file_path)
            img_np = np.array(img)
            
            reader = easyocr.Reader(['en'], gpu=False, verbose=False)
            result = reader.readtext(img_np, detail=0, paragraph=False)
            
            if not result:
                return "No text detected in image"
            
            text = "\n".join(result)
            logger.info(f"✅ Extracted {len(result)} lines from {file_path.name}")
            
            return text
            
        except ImportError:
            logger.error("easyocr not installed")
            raise ImportError("easyocr required. Install: pip install easyocr")
    
    def _convert_text(self, file_path: Path, callback: Optional[Callable] = None) -> str:
        if callback:
            callback(f"Reading text file...")
        
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                text = f.read()
            
            logger.info(f"✅ Read {file_path.name} ({len(text)} chars)")
            return text
            
        except Exception as e:
            logger.error(f"Failed to read {file_path.name}: {str(e)}")
            raise
    
    def get_supported_extensions(self) -> Dict[str, list]:
        return {
            'audio': list(self.audio_extensions),
            'video': list(self.video_extensions),
            'pdf': list(self.pdf_extensions),
            'pptx': list(self.pptx_extensions),
            'docx': list(self.docx_extensions),
            'excel': list(self.excel_extensions),
            'code': list(self.code_extensions),
            'image': list(self.image_extensions),
            'text': list(self.text_extensions)
        }
    
    def is_supported(self, file_path: str) -> bool:
        ext = Path(file_path).suffix.lower()
        all_extensions = (
            self.audio_extensions | self.video_extensions | self.pdf_extensions |
            self.pptx_extensions | self.docx_extensions | self.excel_extensions |
            self.code_extensions | self.image_extensions | self.text_extensions
        )
        return ext in all_extensions
